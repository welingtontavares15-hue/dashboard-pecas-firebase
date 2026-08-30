#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gera um arquivo .pptx a partir de um JSON estruturado criado pela skill do ChatGPT.
Uso:
    python python-generator.py exemplo.json apresentacao.pptx
"""
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Instale a dependência: pip install python-pptx")
    sys.exit(1)


def hex_to_rgb(hex_color: str) -> RgbColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RgbColor(0x33, 0x33, 0x33)
    return RgbColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_text_frame(text_frame, text, font_size, font_name, color, bold=False, align=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    font = p.font
    font.size = Pt(font_size)
    font.name = font_name
    font.bold = bold
    font.color.rgb = color
    return p


def add_title_slide(prs, data, palette, fonts):
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(1.2))
    set_text_frame(title_box.text_frame, data.get("title", ""), 40, fonts["title"], palette["primary"], bold=True, align=PP_ALIGN.CENTER)

    if data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.5), Inches(8.5), Inches(0.8))
        set_text_frame(sub_box.text_frame, data["subtitle"], 22, fonts["body"], palette["text"], align=PP_ALIGN.CENTER)

    if data.get("author"):
        author_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(8.5), Inches(0.5))
        set_text_frame(author_box.text_frame, data["author"], 16, fonts["body"], palette["text"], align=PP_ALIGN.CENTER)

    return slide


def add_content_slide(prs, slide_data, palette, fonts):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.9), Inches(0.9))
    set_text_frame(title_box.text_frame, slide_data.get("title", ""), 32, fonts["title"], palette["primary"], bold=True)

    # Bullets
    bullets = slide_data.get("bullets", [])
    content_text = "\n".join(f"• {b}" for b in bullets)
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    set_text_frame(tf, content_text, 18, fonts["body"], palette["text"])

    # Visual suggestion footer
    visual = slide_data.get("visual", "")
    if visual:
        visual_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.9), Inches(0.6))
        set_text_frame(visual_box.text_frame, f"💡 Visual: {visual}", 12, fonts["body"], palette["secondary"], italic=True)

    # Speaker note
    note = slide_data.get("speaker_note", "")
    if note:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = f"Nota do apresentador:\n{note}"

    return slide


def build_presentation(json_path: str, output_path: str):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    palette_input = data.get("palette", {})
    palette = {
        "primary": hex_to_rgb(palette_input.get("primary", "#0F2C59")),
        "secondary": hex_to_rgb(palette_input.get("secondary", "#3A5A8A")),
        "accent": hex_to_rgb(palette_input.get("accent", "#FF6600")),
        "background": hex_to_rgb(palette_input.get("background", "#FFFFFF")),
        "text": hex_to_rgb(palette_input.get("text", "#333333"))
    }

    fonts_input = data.get("fonts", {})
    fonts = {
        "title": fonts_input.get("title", "Arial"),
        "body": fonts_input.get("body", "Calibri")
    }

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, data, palette, fonts)

    for slide_data in sorted(data.get("slides", []), key=lambda s: s.get("number", 0)):
        add_content_slide(prs, slide_data, palette, fonts)

    prs.save(output_path)
    print(f"Apresentação salva em: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Uso: python python-generator.py <entrada.json> <saida.pptx>")
        sys.exit(1)
    build_presentation(sys.argv[1], sys.argv[2])
